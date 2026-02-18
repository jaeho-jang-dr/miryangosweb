#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
노트랑 PDF → PPTX 변환 모듈

PyMuPDF로 PDF 페이지를 이미지로 추출하고
python-pptx로 PPTX 슬라이드에 삽입합니다.
"""
import logging
import sys
from pathlib import Path
from typing import Optional, Union

if sys.platform == 'win32':
    sys.stdout.reconfigure(encoding='utf-8')

logger = logging.getLogger(__name__)

# Slide dimensions for 16:9 widescreen output
_SLIDE_WIDTH_INCHES: float = 13.333
_SLIDE_HEIGHT_INCHES: float = 7.5
_DEFAULT_DPI: int = 200
# Blank slide layout index in python-pptx default layouts
_BLANK_LAYOUT_INDEX: int = 6


def pdf_to_pptx(
    pdf_path: Union[str, Path],
    pptx_path: Optional[Union[str, Path]] = None,
    dpi: int = _DEFAULT_DPI,
) -> bool:
    """Convert a PDF file to a PPTX presentation.

    Each PDF page is rasterised at *dpi* and embedded as a full-slide image
    in the resulting PPTX file.

    Args:
        pdf_path: Path to the source PDF file (str or Path).
        pptx_path: Destination PPTX path. Defaults to the same stem as
            *pdf_path* with a ``.pptx`` extension.
        dpi: Rasterisation resolution in dots-per-inch (default 200).

    Returns:
        ``True`` on success, ``False`` otherwise.
    """
    try:
        import fitz  # PyMuPDF
        from pptx import Presentation
        from pptx.util import Inches
        from io import BytesIO
    except ImportError as e:
        logger.error("Missing required package: %s", e)
        print(f"필요한 패키지가 설치되지 않았습니다: {e}")
        print("설치: pip install pymupdf python-pptx")
        return False

    pdf_path = Path(pdf_path)
    if not pdf_path.exists():
        logger.error("PDF file not found: %s", pdf_path)
        print(f"PDF 파일을 찾을 수 없습니다: {pdf_path}")
        return False

    pptx_path = Path(pptx_path) if pptx_path is not None else pdf_path.with_suffix('.pptx')

    try:
        pdf_doc = fitz.open(str(pdf_path))
        page_count = len(pdf_doc)

        logger.info("Converting %s (%d pages)", pdf_path.name, page_count)
        print(f"PDF 변환 중: {pdf_path.name} ({page_count}페이지)")

        prs = Presentation()
        prs.slide_width = Inches(_SLIDE_WIDTH_INCHES)
        prs.slide_height = Inches(_SLIDE_HEIGHT_INCHES)

        blank_layout = prs.slide_layouts[_BLANK_LAYOUT_INDEX]

        # Build fitz render matrix once; cache slide dimensions to avoid repeated attribute access
        mat = fitz.Matrix(dpi / 72, dpi / 72)
        slide_w = prs.slide_width
        slide_h = prs.slide_height

        for page_num in range(page_count):
            page = pdf_doc[page_num]
            pix = page.get_pixmap(matrix=mat)
            img_bytes = BytesIO(pix.tobytes("png"))

            slide = prs.slides.add_slide(blank_layout)
            slide.shapes.add_picture(
                img_bytes,
                Inches(0),
                Inches(0),
                width=slide_w,
                height=slide_h,
            )

            img_bytes.close()
            pix = None

            print(f"  페이지 {page_num + 1}/{page_count} 변환 완료")

        pdf_doc.close()

        prs.save(str(pptx_path))
        print(f"PPTX 저장 완료: {pptx_path}")

        return True

    except Exception as e:
        print(f"변환 실패: {e}")
        return False


def batch_convert(
    pdf_dir: Union[str, Path],
    output_dir: Optional[Union[str, Path]] = None,
    dpi: int = _DEFAULT_DPI,
) -> int:
    """Convert all PDF files found in *pdf_dir* to PPTX.

    Args:
        pdf_dir: Directory containing source PDF files.
        output_dir: Destination directory for PPTX files. Defaults to *pdf_dir*.
        dpi: Rasterisation resolution in dots-per-inch (default 200).

    Returns:
        Number of files converted successfully.
    """
    pdf_dir = Path(pdf_dir)
    if not pdf_dir.exists():
        print(f"디렉토리를 찾을 수 없습니다: {pdf_dir}")
        return 0

    if output_dir:
        output_dir = Path(output_dir)
        output_dir.mkdir(parents=True, exist_ok=True)
    else:
        output_dir = pdf_dir

    pdf_files = list(pdf_dir.glob("*.pdf"))
    if not pdf_files:
        print(f"PDF 파일이 없습니다: {pdf_dir}")
        return 0

    print(f"\n{len(pdf_files)}개 PDF 파일 변환 시작\n")

    success_count = 0
    for pdf_file in pdf_files:
        pptx_path = output_dir / pdf_file.with_suffix('.pptx').name

        if pdf_to_pptx(str(pdf_file), str(pptx_path), dpi):
            success_count += 1
        print()

    print(f"\n변환 완료: {success_count}/{len(pdf_files)} 파일")
    return success_count


# CLI 실행
if __name__ == "__main__":
    import argparse

    parser = argparse.ArgumentParser(description="PDF → PPTX 변환")
    parser.add_argument("pdf_path", help="PDF 파일 또는 디렉토리 경로")
    parser.add_argument("-o", "--output", help="출력 경로")
    parser.add_argument("--dpi", type=int, default=200, help="이미지 해상도 (기본 200)")
    parser.add_argument("--batch", action="store_true", help="디렉토리 일괄 변환")

    args = parser.parse_args()

    if args.batch:
        batch_convert(args.pdf_path, args.output, args.dpi)
    else:
        success = pdf_to_pptx(args.pdf_path, args.output, args.dpi)
        sys.exit(0 if success else 1)
