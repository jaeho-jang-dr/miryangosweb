#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
노트랑 변환 모듈
- PDF → PPTX 변환
- 슬라이드에 노트 추가
"""
import io
import logging
import sys
from pathlib import Path
from typing import Optional, Tuple, List, Dict

if sys.platform == 'win32':
    sys.stdout.reconfigure(encoding='utf-8')

logger = logging.getLogger(__name__)

# Default slide dimensions (16:9 widescreen)
_SLIDE_WIDTH_INCHES: float = 13.333
_SLIDE_HEIGHT_INCHES: float = 7.5
_DEFAULT_ZOOM: float = 2.0
_DEFAULT_TITLE_FONT_SIZE: int = 44
_DEFAULT_BODY_FONT_SIZE: int = 24
# Maximum characters kept from PDF pages for title and body
_MAX_TITLE_CHARS: int = 100
_MAX_BODY_CHARS: int = 1000


def pdf_to_pptx(
    pdf_path: Path,
    output_path: Optional[Path] = None,
    zoom: float = _DEFAULT_ZOOM,
    slide_width_inches: float = _SLIDE_WIDTH_INCHES,
    slide_height_inches: float = _SLIDE_HEIGHT_INCHES,
) -> Tuple[Path, int]:
    """
    PDF를 PPTX로 변환

    Args:
        pdf_path: PDF 파일 경로
        output_path: 출력 경로 (None이면 같은 이름으로)
        zoom: 이미지 확대 비율 (2.0 = 2배)
        slide_width_inches: 슬라이드 너비 (인치)
        slide_height_inches: 슬라이드 높이 (인치)

    Returns:
        (출력 파일 경로, 슬라이드 수)
    """
    import fitz
    from pptx import Presentation
    from pptx.util import Inches

    pdf_path = Path(pdf_path)
    out_path = Path(output_path) if output_path else pdf_path.with_suffix('.pptx')

    # PERF: 단일 Matrix 객체 재사용 (루프마다 재생성 방지)
    mat = fitz.Matrix(zoom, zoom)

    # PDF 열기
    doc = fitz.open(pdf_path)

    # 새 프레젠테이션 생성
    prs = Presentation()
    prs.slide_width = Inches(slide_width_inches)
    prs.slide_height = Inches(slide_height_inches)

    # PERF: 슬라이드 크기를 지역 변수로 캐시 (반복 속성 조회 방지)
    slide_w = prs.slide_width
    slide_h = prs.slide_height

    # 빈 레이아웃 사용
    blank_layout = prs.slide_layouts[6]  # 완전 빈 레이아웃

    # 각 페이지를 슬라이드로 변환
    for page_num in range(len(doc)):
        page = doc[page_num]

        # PDF 페이지를 이미지로 렌더링
        pix = page.get_pixmap(matrix=mat)
        # PERF: PNG 대신 JPEG 사용 시 파일 크기/메모리 절감 가능하나
        # 품질 유지를 위해 PNG 유지. tobytes()는 in-memory이므로 직접 BytesIO로 전달
        img_data = pix.tobytes("png")

        # 슬라이드 추가
        slide = prs.slides.add_slide(blank_layout)

        # 이미지를 전체 슬라이드로 추가
        slide.shapes.add_picture(
            io.BytesIO(img_data),
            Inches(0), Inches(0),
            width=slide_w,
            height=slide_h
        )

        # PERF: 각 페이지 렌더링 후 픽셀맵 명시 해제로 메모리 조기 반환
        del img_data
        pix = None

    doc.close()
    prs.save(out_path)

    return out_path, len(prs.slides)


def add_notes_to_pptx(
    pptx_path: Path,
    notes: List[str],
    output_path: Optional[Path] = None,
) -> Path:
    """
    PPTX 슬라이드에 노트 추가

    Args:
        pptx_path: PPTX 파일 경로
        notes: 각 슬라이드의 노트 목록
        output_path: 출력 경로 (None이면 덮어쓰기)

    Returns:
        출력 파일 경로
    """
    from pptx import Presentation

    pptx_path = Path(pptx_path)
    out_path = Path(output_path) if output_path else pptx_path

    prs = Presentation(pptx_path)

    for i, slide in enumerate(prs.slides):
        if i < len(notes) and notes[i]:
            # 노트 슬라이드 가져오기 또는 생성
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = notes[i]

    prs.save(out_path)
    return out_path


def extract_text_from_pdf(pdf_path: Path) -> List[str]:
    """
    PDF에서 텍스트 추출 (페이지별)

    Args:
        pdf_path: PDF 파일 경로

    Returns:
        페이지별 텍스트 목록
    """
    import fitz

    doc = fitz.open(pdf_path)
    # PERF: 리스트 컴프리헨션으로 루프 오버헤드 절감
    texts = [page.get_text() for page in doc]
    doc.close()
    return texts


def pdf_to_pptx_with_notes(
    pdf_path: Path,
    output_path: Optional[Path] = None,
    custom_notes: Optional[List[str]] = None,
) -> Tuple[Path, int]:
    """
    PDF를 PPTX로 변환하고 노트 추가

    Args:
        pdf_path: PDF 파일 경로
        output_path: 출력 경로
        custom_notes: 사용자 정의 노트 (None이면 PDF 텍스트 사용)

    Returns:
        (출력 파일 경로, 슬라이드 수)
    """
    # PDF를 PPTX로 변환
    pptx_path, slide_count = pdf_to_pptx(pdf_path, output_path)

    # 노트 추가
    if custom_notes:
        add_notes_to_pptx(pptx_path, custom_notes)
    else:
        # PDF 텍스트를 노트로 추가
        texts = extract_text_from_pdf(pdf_path)
        add_notes_to_pptx(pptx_path, texts)

    return pptx_path, slide_count


def batch_convert(
    pdf_dir: Path,
    output_dir: Optional[Path] = None,
    pattern: str = "*.pdf",
) -> List[Tuple[Path, Optional[Path], int]]:
    """
    여러 PDF를 일괄 변환

    Args:
        pdf_dir: PDF 디렉토리
        output_dir: 출력 디렉토리 (None이면 같은 디렉토리)
        pattern: 파일 패턴

    Returns:
        [(원본, 출력, 슬라이드수), ...]
    """
    pdf_dir = Path(pdf_dir)
    out_dir = Path(output_dir) if output_dir else pdf_dir
    out_dir.mkdir(parents=True, exist_ok=True)

    results = []

    for pdf_file in pdf_dir.glob(pattern):
        out_file = out_dir / pdf_file.with_suffix('.pptx').name

        try:
            pptx_path, count = pdf_to_pptx(pdf_file, out_file)
            results.append((pdf_file, pptx_path, count))
            print(f"  {pdf_file.name} → {pptx_path.name} ({count} slides)")
        except Exception as e:
            logger.error("Conversion failed for %s: %s", pdf_file.name, e)
            print(f"  {pdf_file.name} 변환 실패: {e}")
            results.append((pdf_file, None, 0))

    return results


def apply_template(
    pptx_path: Path,
    template_path: Path,
    output_path: Optional[Path] = None,
) -> Path:
    """
    PPTX에 템플릿 스타일 적용

    Args:
        pptx_path: 원본 PPTX 파일
        template_path: 템플릿 PPTX 파일
        output_path: 출력 경로 (None이면 덮어쓰기)

    Returns:
        출력 파일 경로
    """
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from copy import deepcopy

    pptx_path = Path(pptx_path)
    template_path = Path(template_path)
    out_path = Path(output_path) if output_path else pptx_path

    # 템플릿 로드
    template = Presentation(template_path)

    # 원본 로드
    prs = Presentation(pptx_path)

    # 템플릿에서 슬라이드 크기 복사
    prs.slide_width = template.slide_width
    prs.slide_height = template.slide_height

    # 템플릿 배경/테마 정보 추출 (있는 경우)
    if len(template.slides) > 0:
        template_slide = template.slides[0]
        # 배경색 복사 시도
        try:
            if hasattr(template_slide, 'background'):
                for slide in prs.slides:
                    slide.background.fill.solid()
                    if template_slide.background.fill.type:
                        slide.background.fill._fill = deepcopy(template_slide.background.fill._fill)
        except Exception:
            logger.debug("Background copy from template failed; ignoring.", exc_info=True)

    prs.save(out_path)
    return out_path


def create_styled_pptx(
    content_data: List[Dict],
    output_path: Path,
    style: str = "modern",
    title_font_size: int = _DEFAULT_TITLE_FONT_SIZE,
    body_font_size: int = _DEFAULT_BODY_FONT_SIZE,
) -> Path:
    """
    스타일이 적용된 PPTX 생성

    Args:
        content_data: [{"title": "제목", "body": "내용", "image": "이미지경로"}, ...]
        output_path: 출력 경로
        style: 스타일 ("modern", "minimal", "corporate", "creative")
        title_font_size: 제목 폰트 크기 (pt)
        body_font_size: 본문 폰트 크기 (pt)

    Returns:
        출력 파일 경로
    """
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    # PERF: MSO_SHAPE을 루프 밖에서 한 번만 import (기존 코드는 루프 내부에서 반복 import)
    from pptx.enum.shapes import MSO_SHAPE

    # 스타일별 색상 테마
    STYLES = {
        "modern": {
            "bg": RGBColor(0x1a, 0x1a, 0x2e),  # 다크 네이비
            "title": RGBColor(0xff, 0xff, 0xff),  # 화이트
            "body": RGBColor(0xe0, 0xe0, 0xe0),  # 라이트 그레이
            "accent": RGBColor(0x00, 0xd9, 0xff),  # 사이안
        },
        "minimal": {
            "bg": RGBColor(0xff, 0xff, 0xff),  # 화이트
            "title": RGBColor(0x33, 0x33, 0x33),  # 다크 그레이
            "body": RGBColor(0x66, 0x66, 0x66),  # 미디엄 그레이
            "accent": RGBColor(0x00, 0x7a, 0xcc),  # 블루
        },
        "corporate": {
            "bg": RGBColor(0xf5, 0xf5, 0xf5),  # 라이트 그레이
            "title": RGBColor(0x00, 0x3d, 0x7a),  # 다크 블루
            "body": RGBColor(0x33, 0x33, 0x33),  # 다크 그레이
            "accent": RGBColor(0x00, 0x7a, 0xcc),  # 미디엄 블루
        },
        "creative": {
            "bg": RGBColor(0x2d, 0x13, 0x2c),  # 다크 퍼플
            "title": RGBColor(0xff, 0xd9, 0x3d),  # 골드
            "body": RGBColor(0xff, 0xff, 0xff),  # 화이트
            "accent": RGBColor(0xff, 0x6b, 0x6b),  # 코랄
        },
    }

    colors = STYLES.get(style, STYLES["modern"])

    prs = Presentation()
    prs.slide_width = Inches(_SLIDE_WIDTH_INCHES)
    prs.slide_height = Inches(_SLIDE_HEIGHT_INCHES)

    blank_layout = prs.slide_layouts[6]

    for item in content_data:
        slide = prs.slides.add_slide(blank_layout)

        # 배경색 설정
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = colors["bg"]

        # 제목 텍스트박스
        if "title" in item and item["title"]:
            title_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(0.5),
                Inches(12.333), Inches(1.2)
            )
            title_frame = title_box.text_frame
            title_frame.word_wrap = True
            title_p = title_frame.paragraphs[0]
            title_p.text = item["title"]
            title_p.font.size = Pt(title_font_size)
            title_p.font.bold = True
            title_p.font.color.rgb = colors["title"]
            title_p.alignment = PP_ALIGN.LEFT

        # 본문 텍스트박스
        if "body" in item and item["body"]:
            body_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(2.0),
                Inches(12.333), Inches(5.0)
            )
            body_frame = body_box.text_frame
            body_frame.word_wrap = True
            body_p = body_frame.paragraphs[0]
            body_p.text = item["body"]
            body_p.font.size = Pt(body_font_size)
            body_p.font.color.rgb = colors["body"]
            body_p.alignment = PP_ALIGN.LEFT

        # 이미지 (있는 경우)
        if "image" in item and item["image"] and Path(item["image"]).exists():
            slide.shapes.add_picture(
                item["image"],
                Inches(8.0), Inches(2.0),
                width=Inches(4.833)
            )

        # 악센트 라인
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.5), Inches(1.8),
            Inches(2.0), Inches(0.05)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = colors["accent"]
        line.line.fill.background()

    prs.save(output_path)
    return output_path


class Converter:
    """변환기 클래스"""

    def __init__(self, output_dir: Path = None):
        from .config import get_config
        config = get_config()
        self.output_dir = Path(output_dir) if output_dir else config.download_dir
        self.templates_dir = Path(__file__).parent / "templates"

    def pdf_to_pptx(self, pdf_path: Path) -> Tuple[Optional[Path], int]:
        """Convert a PDF to PPTX, placing the output in :attr:`output_dir`.

        Args:
            pdf_path: Source PDF file path.

        Returns:
            Tuple of (output PPTX path, slide count). Returns ``(None, 0)`` on error.
        """
        try:
            output_path = self.output_dir / Path(pdf_path).with_suffix('.pptx').name
            return pdf_to_pptx(pdf_path, output_path)
        except Exception as e:
            logger.error("PDF to PPTX conversion failed: %s", e)
            print(f"  변환 실패: {e}")
            return None, 0

    def pdf_to_pptx_with_notes(
        self,
        pdf_path: Path,
        notes: Optional[List[str]] = None,
    ) -> Tuple[Optional[Path], int]:
        """Convert a PDF to PPTX and embed speaker notes.

        Args:
            pdf_path: Source PDF file path.
            notes: Per-slide note strings. Uses extracted PDF text when ``None``.

        Returns:
            Tuple of (output PPTX path, slide count). Returns ``(None, 0)`` on error.
        """
        try:
            output_path = self.output_dir / Path(pdf_path).with_suffix('.pptx').name
            return pdf_to_pptx_with_notes(pdf_path, output_path, notes)
        except Exception as e:
            logger.error("PDF to PPTX (with notes) conversion failed: %s", e)
            print(f"  변환 실패: {e}")
            return None, 0

    def batch_convert(
        self,
        pdf_dir: Path,
        pattern: str = "*.pdf",
    ) -> List[Tuple[Path, Optional[Path], int]]:
        """Batch-convert all PDFs matching *pattern* in *pdf_dir*.

        Args:
            pdf_dir: Directory containing source PDF files.
            pattern: Glob pattern to match files (default ``"*.pdf"``).

        Returns:
            List of ``(source_pdf, output_pptx, slide_count)`` tuples.
        """
        return batch_convert(pdf_dir, self.output_dir, pattern)

    def pdf_to_styled_pptx(
        self,
        pdf_path: Path,
        style: str = "modern",
        output_path: Optional[Path] = None,
    ) -> Tuple[Optional[Path], int]:
        """Convert a PDF to a styled PPTX by extracting text from each page.

        Args:
            pdf_path: Source PDF file path.
            style: Visual style preset (``"modern"``, ``"minimal"``,
                ``"corporate"``, or ``"creative"``).
            output_path: Destination PPTX path. Auto-generated when ``None``.

        Returns:
            Tuple of (output PPTX path, slide count). Returns ``(None, 0)`` on error.
        """
        try:
            texts = extract_text_from_pdf(pdf_path)

            content_data: List[Dict] = []
            for i, text in enumerate(texts):
                lines = text.strip().split('\n')
                title = lines[0] if lines else f"슬라이드 {i + 1}"
                body = '\n'.join(lines[1:]) if len(lines) > 1 else ""
                content_data.append({
                    "title": title[:_MAX_TITLE_CHARS],
                    "body": body[:_MAX_BODY_CHARS],
                })

            if output_path is None:
                output_path = self.output_dir / f"{Path(pdf_path).stem}_styled.pptx"

            result = create_styled_pptx(content_data, output_path, style=style)
            return result, len(content_data)
        except Exception as e:
            logger.error("Styled PPTX conversion failed: %s", e)
            print(f"  스타일 변환 실패: {e}")
            return None, 0

    def apply_template(
        self,
        pptx_path: Path,
        template_path: Path,
        output_path: Optional[Path] = None,
    ) -> Optional[Path]:
        """Apply a PPTX template to an existing presentation.

        Args:
            pptx_path: Source PPTX file path.
            template_path: Template PPTX file path.
            output_path: Destination path. Auto-generated when ``None``.

        Returns:
            Path to the updated PPTX file, or ``None`` on failure.
        """
        try:
            if output_path is None:
                output_path = self.output_dir / f"{Path(pptx_path).stem}_templated.pptx"
            return apply_template(pptx_path, template_path, output_path)
        except Exception as e:
            logger.error("Template application failed: %s", e)
            print(f"  템플릿 적용 실패: {e}")
            return None

    def create_from_data(
        self,
        content_data: List[Dict],
        title: str,
        style: str = "modern",
    ) -> Optional[Path]:
        """Generate a styled PPTX from structured content data.

        Args:
            content_data: List of slide content dicts with ``"title"`` and
                optional ``"body"`` keys.
            title: Base filename (without extension) for the output file.
            style: Visual style preset (``"modern"``, ``"minimal"``,
                ``"corporate"``, or ``"creative"``).

        Returns:
            Path to the generated PPTX file, or ``None`` on failure.
        """
        try:
            output_path = self.output_dir / f"{title}.pptx"
            return create_styled_pptx(content_data, output_path, style=style)
        except Exception as e:
            logger.error("PPTX creation from data failed: %s", e)
            print(f"  생성 실패: {e}")
            return None
